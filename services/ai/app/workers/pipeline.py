"""Document ingestion pipeline — QVAC-primary, structure-aware chunking.

Flow: parse (page-by-page) → clean → chunk (structure-aware) → filter → JSONL → QVAC /ingest
ChromaDB is not written during ingestion. It remains as a passive fallback
at query time in chat_service.py if QVAC is unreachable.
"""
import json
import logging
import os
import re
from collections import Counter
from pathlib import Path

import httpx

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Path and env constants
# ---------------------------------------------------------------------------
_HERE = Path(__file__).resolve()
_SERVICES_AI = _HERE.parents[2]          # services/ai/

CHROMA_DB_PATH = os.getenv("CHROMA_DB_PATH", str(_SERVICES_AI / "chroma_db"))
CHROMA_COLLECTION_NAME = os.getenv("CHROMA_COLLECTION_NAME", "bitpolito_course")
UPLOADS_DIR = Path(os.getenv("UPLOADS_DIR", str(_SERVICES_AI / "uploads")))
QVAC_INGEST_DIR = Path(os.getenv("QVAC_INGEST_DIR", str(_SERVICES_AI / "qvac_ingest")))
QVAC_SERVICE_URL = os.getenv("QVAC_SERVICE_URL", "http://localhost:3001")
SKIP_CHROMA_INDEX = os.getenv("SKIP_CHROMA_INDEX", "false").lower() == "true"

from app.db.models import CourseDocument, DocumentProcessingStage, DocumentStatus  # noqa: E402
from app.db.session import get_db_context                                           # noqa: E402
from app.repositories import document_repo                                          # noqa: E402

# ---------------------------------------------------------------------------
# sys.modules aliasing — dual-import guard
# ---------------------------------------------------------------------------
# Register 'services.ai.app.*' as aliases for 'app.*' in sys.modules.
# This ensures that classes imported via either path are the same objects,
# preventing silent Pydantic isinstance failures when the worker is invoked
# from the project root instead of from services/ai/.
import sys as _sys
import types as _types


def _register_module_aliases() -> None:
    canonical = "app"
    alias_root = "services.ai.app"

    for ns_name in ("services", "services.ai", "services.ai.app"):
        if ns_name not in _sys.modules:
            ns = _types.ModuleType(ns_name)
            ns.__path__ = []  # type: ignore[attr-defined]
            _sys.modules[ns_name] = ns

    for name in list(_sys.modules):
        if name == canonical or name.startswith(canonical + "."):
            long_name = alias_root + name[len(canonical):]
            _sys.modules.setdefault(long_name, _sys.modules[name])


_register_module_aliases()

# ---------------------------------------------------------------------------
# Chunking parameters
# ---------------------------------------------------------------------------
_PARENT_WORDS = 1200    # parent chunk: LLM context window (≈ 1500 tokens)
_CHILD_WORDS = 150      # child chunk: retrieval unit (≈ 200 tokens)
_CHILD_MAX_WORDS = 350  # hard cap: single long sentence can exceed 150-word target; 350 ≈ 455 tokens (GTE-Large limit 512)
_CHILD_OVERLAP = 30     # overlap between consecutive child chunks (words)
_MAX_WORDS = 400        # legacy: only used by chunk_pages() (no longer called by run())
_OVERLAP_WORDS = 50     # legacy: overlap used by chunk_pages()
_MIN_WORDS = 25         # paragraph threshold: shorter chunks are discarded
_MIN_WORDS_TABLE = 4    # table threshold: one data row is enough (cells are short)

_RAG_CONTEXTUAL_CHUNKS = os.getenv("RAG_CONTEXTUAL_CHUNKS", "false").lower() == "true"
_CONTEXTUAL_TIMEOUT = float(os.getenv("QVAC_CONTEXTUAL_TIMEOUT", "25"))

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _word_count(text: str) -> int:
    return len(text.split())


_FIGURE_CAPTION_RE = re.compile(r'_Figure\s+[\d][\d.-]*[^_]*_', re.IGNORECASE)

# Patterns da rimuovere da ogni pagina
_STRIP_PATTERNS = [
    re.compile(r'www\.\S+\.ir\b[^\S\n]*', re.IGNORECASE),       # watermark EBooksWorld e simili
    re.compile(r'www\.\S+\.com/?\s*\n', re.IGNORECASE),          # altri watermark URL inline
    re.compile(r'^\s*\d+\s*\|\s*Chapter[^\n]*', re.MULTILINE),   # "8 | Chapter 1: Introduction"
    re.compile(r'[­​‌‍﻿]'),                                         # unicode invisibili (soft-hyphen, ZWS, ecc.)
    # PDF running-header artifacts produced by pymupdf4llm for LaTeX books:
    # "## 32"  — even-page running header (just a page number formatted as heading)
    # "## 34 T H E B I T C O I N S T A N D A R D" — odd-page: page# + spaced book title
    re.compile(r'^#{1,4}\s+\d+(?:\s+(?:[A-Z]\s+)+[A-Z]+)?\s*$', re.MULTILINE),
    # Standalone book page numbers in paragraph text (e.g. "219" alone on a line)
    re.compile(r'^\d{1,4}\s*$', re.MULTILINE),
    # LaTeX source metadata lines injected by pymupdf4llm (typesetting artefacts):
    # "Ammous c01.tex V1 - 03/05/2018 1:08pm Page 10"
    re.compile(r'^[A-Za-z]+\s+\w+\.tex\s+V\d+[^\n]*$', re.MULTILINE),
]

# Regex per sentence splitting (fine frase + inizio maiuscola)
_SENTENCE_SPLIT_RE = re.compile(r'(?<=[.!?])\s+(?=[A-Z\"])')


# ---------------------------------------------------------------------------
# Stage 0 — Docling PDF parser (optional, activated by USE_DOCLING=true)
# ---------------------------------------------------------------------------

def _parse_pdf_with_docling(file_path: str) -> tuple[list[dict], int]:
    """Docling-based PDF parser — returns the same pages format as parse_pdf_pages.

    Docling produces higher-quality structured extraction (better table handling,
    heading detection, formula recognition) than pymupdf4llm. Activated when
    USE_DOCLING=true; otherwise parse_pdf_pages (pymupdf4llm) is used.

    Falls back to parse_pdf_pages if Docling is not installed or conversion fails.
    """
    from collections import defaultdict
    from docling.document_converter import DocumentConverter

    converter = DocumentConverter()
    result = converter.convert(file_path)
    doc = result.document

    page_texts: dict[int, list[str]] = defaultdict(list)

    for item, _ in doc.iterate_items():
        prov_list = getattr(item, "prov", None)
        prov = prov_list[0] if prov_list else None
        page_no = int(prov.page_no) if prov else 0

        # Tables: export to markdown for structured representation
        raw_text: str = ""
        try:
            from docling_core.types.doc.document import TableItem
            if isinstance(item, TableItem):
                try:
                    raw_text = item.export_to_dataframe().to_markdown(index=False) or ""
                except Exception:
                    raw_text = getattr(item, "text", None) or ""
            else:
                raw_text = getattr(item, "text", None) or ""
        except ImportError:
            raw_text = getattr(item, "text", None) or ""

        text = raw_text.strip()
        if not text:
            continue
        # Fix PDF ligature corruption (same as StructuralParser._sanitize_text)
        text = text.replace("昀椀", "fi").replace("昀氀", "fl")
        page_texts[page_no].append(text)

    pages = [
        {"page": pno, "text": "\n\n".join(blocks)}
        for pno, blocks in sorted(page_texts.items())
        if any(b.strip() for b in blocks)
    ]
    page_count = len(getattr(doc, "pages", pages))
    return pages, page_count


_USE_DOCLING = os.getenv("USE_DOCLING", "false").lower() == "true"


# ---------------------------------------------------------------------------
# Stage 1 — Parsers (page-by-page)
# ---------------------------------------------------------------------------

def parse_pdf_pages(file_path: str) -> tuple[list[dict], int]:
    """pymupdf4llm page_chunks=True → [{page, text}] + page count.

    Ogni elemento corrisponde a una pagina PDF con il suo numero (1-indexed).
    """
    import pymupdf4llm
    import fitz

    raw = pymupdf4llm.to_markdown(file_path, page_chunks=True)
    pages = []
    for p in raw:
        meta = p.get("metadata", {}) if isinstance(p, dict) else {}
        # get_metadata in pymupdf4llm sets page = pno + 1 (already 1-indexed)
        # pymupdf4llm ≥1.27 uses "page_number" (1-indexed); older builds used "page"
        page_num = meta.get("page_number") or meta.get("page") or 0
        text = p.get("text", "") if isinstance(p, dict) else str(p)
        pages.append({"page": page_num, "text": text})

    with fitz.open(file_path) as pdf:
        page_count = pdf.page_count

    return pages, page_count


def parse_pptx_pages(file_path: str) -> tuple[list[dict], int]:
    """python-pptx → [{page=slide_num, text}] per ogni slide."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = (
            slide.shapes.title.text.strip()
            if slide.shapes.title and slide.shapes.title.has_text_frame
            else f"Slide {i}"
        )
        body_parts = []
        for shape in slide.placeholders:
            if (
                shape.has_text_frame
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx != 0
            ):
                body_parts.append(shape.text_frame.text.strip())  # type: ignore[union-attr]
        body = "\n".join(p for p in body_parts if p)
        notes = ""
        if slide.has_notes_slide:
            nf = slide.notes_slide.notes_text_frame
            notes = nf.text.strip() if nf else ""
        parts = [f"## {title}"]
        if body:
            parts.append(body)
        if notes:
            parts.append(f"*Notes: {notes}*")
        slides.append({"page": i, "text": "\n\n".join(parts)})

    return slides, len(prs.slides)


def parse_docx_pages(file_path: str) -> tuple[list[dict], int]:
    """python-docx → singola entry [{page=1, text}] (nessun page count affidabile)."""
    from docx import Document

    doc = Document(file_path)
    lines = []
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        style_name = (para.style.name or "") if para.style is not None else ""
        if style_name.startswith("Heading"):
            try:
                level = int(style_name.split()[-1])
            except ValueError:
                level = 2
            lines.append(f"{'#' * level} {para.text.strip()}")
        else:
            lines.append(para.text.strip())

    return [{"page": 1, "text": "\n\n".join(lines)}], 0


# ---------------------------------------------------------------------------
# Stage 2 — Cleaner
# ---------------------------------------------------------------------------

def detect_boilerplate(pages: list[dict], min_freq: float = 0.3) -> set[str]:
    """Identifica righe che compaiono in >= min_freq delle pagine (header/footer).

    Ritorna un set di stringhe stripped da rimuovere durante il cleaning.
    Ignorato per documenti con < 5 pagine (troppo poco campione).
    """
    if len(pages) < 5:
        return set()

    line_counts: Counter = Counter()
    for p in pages:
        seen_on_page: set[str] = set()
        for line in p["text"].splitlines():
            s = line.strip()
            if s and len(s) > 3:
                seen_on_page.add(s)
        line_counts.update(seen_on_page)

    threshold = len(pages) * min_freq
    return {line for line, count in line_counts.items() if count >= threshold}


def clean_page(text: str, boilerplate: set[str]) -> str:
    """Rimuove watermark, boilerplate ripetuti e unicode invisibili da una pagina."""
    # Rimuovi pattern fissi
    for pat in _STRIP_PATTERNS:
        text = pat.sub("", text)

    # Rimuovi righe di boilerplate rilevate automaticamente
    if boilerplate:
        lines = []
        for line in text.splitlines():
            if line.strip() not in boilerplate:
                lines.append(line)
        text = "\n".join(lines)

    # Comprimi newline multipli
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


# ---------------------------------------------------------------------------
# Stage 3 — Structure-aware chunker
# ---------------------------------------------------------------------------

_LATEX_BLOCK_RE = re.compile(r'\$\$.+?\$\$', re.DOTALL)
_CODE_FENCE_RE = re.compile(r'```.*?```', re.DOTALL)

# Spurious headings emitted by pymupdf4llm from LaTeX running headers / artifacts:
#   - Pure page number:                 "32", "44"
#   - Page + spaced book title:         "34 T H E B I T C O I N S T A N D A R D"
#   - Pure spaced-caps (section name):  "P R O L O G U E", "B I B L I O G R A P H Y"
#   - Roman-numeral + spaced-caps:      "viii C O N T E N T S"
#   - Decorative strikethrough:         "~~k~~" (PDF ornament converted to markdown)
#   - Lowercase Roman numeral labels:   "xiv", "xviii" (front-matter page labels)
_SPURIOUS_HEADING_RE = re.compile(
    r'^\d+(?:\s+(?:[A-Z]\s+)+[A-Z]+)?\s*$'         # digit or digit + spaced-caps
    r'|'
    r'^(?:[a-z]+\s+)?(?:[A-Z]\s+){3,}[A-Z]+\s*$',  # pure/roman-prefix spaced-caps (≥4 caps)
)
_ROMAN_NUMERAL_RE = re.compile(r'^[ivxlcdm]+$')  # lowercase Roman numerals (front-matter labels)


def _is_spurious_heading(heading_text: str) -> bool:
    """Returns True for PDF running-header lines formatted as markdown headings."""
    bare = re.sub(r'^#{1,4}\s+', '', heading_text).strip()
    if _SPURIOUS_HEADING_RE.match(bare):
        return True
    # Strikethrough-only content — decorative PDF artifacts (e.g. "~~k~~")
    if not re.sub(r'~~[^~]*~~', '', bare).strip():
        return True
    # Standalone lowercase Roman numeral (front-matter page labels: xiv, xviii …)
    if _ROMAN_NUMERAL_RE.match(bare):
        return True
    return False


def _split_into_blocks(text: str) -> list[dict]:
    """Segmenta il testo in blocchi tipizzati: heading | table | formula | code | paragraph.

    Formula ($$...$$) e code fence (```...```) vengono estratti come blocchi atomici
    prima del parsing riga per riga, per preservarli intatti.
    """
    # Extract LaTeX block formulas and code fences as standalone blocks first,
    # replacing them with a sentinel so line-by-line parsing is not disrupted.
    sentinel_map: dict[str, dict] = {}

    def _replace(m: re.Match, block_type: str) -> str:
        key = f"\x00BLOCK{len(sentinel_map)}\x00"
        sentinel_map[key] = {"type": block_type, "text": m.group(0).strip()}
        return f"\n{key}\n"

    text = _LATEX_BLOCK_RE.sub(lambda m: _replace(m, "formula"), text)
    text = _CODE_FENCE_RE.sub(lambda m: _replace(m, "code"), text)

    blocks: list[dict] = []
    current_type: str = "paragraph"
    current_lines: list[str] = []

    def flush() -> None:
        if current_lines:
            t = "".join(current_lines).strip()
            if t:
                blocks.append({"type": current_type, "text": t})
            current_lines.clear()

    for line in text.splitlines(keepends=True):
        stripped = line.strip()

        if stripped in sentinel_map:
            flush()
            current_type = "paragraph"
            blocks.append(sentinel_map[stripped])
        elif re.match(r'^#{1,4}\s+\S', stripped):
            if _is_spurious_heading(stripped):
                continue  # discard page-number / running-header fake headings
            flush()
            blocks.append({"type": "heading", "text": stripped})
        elif stripped.startswith("|"):
            if current_type != "table":
                flush()
                current_type = "table"
            current_lines.append(line)
        else:
            if current_type == "table":
                if not stripped:
                    flush()
                    current_type = "paragraph"
                else:
                    flush()
                    current_type = "paragraph"
                    current_lines.append(line)
            else:
                current_type = "paragraph"
                current_lines.append(line)

    flush()
    return blocks


def _split_paragraph(text: str, max_words: int, overlap_words: int = 0) -> list[str]:
    """Divide un paragrafo lungo in sub-chunk ancorati a fine frase.

    overlap_words > 0 aggiunge una sliding window: ogni sub-chunk (tranne il primo)
    inizia con le ultime ~overlap_words parole del chunk precedente.
    """
    if _word_count(text) <= max_words:
        return [text]

    sentences = _SENTENCE_SPLIT_RE.split(text)
    result: list[str] = []
    current: list[str] = []
    current_words = 0

    for sent in sentences:
        sent_words = _word_count(sent)
        if current_words + sent_words > max_words and current:
            result.append(" ".join(current))
            if overlap_words > 0:
                overlap: list[str] = []
                overlap_count = 0
                for s in reversed(current):
                    w = _word_count(s)
                    if overlap_count + w > overlap_words:
                        break
                    overlap.insert(0, s)
                    overlap_count += w
                current = overlap
                current_words = overlap_count
            else:
                current = []
                current_words = 0
        current.append(sent)
        current_words += sent_words

    if current:
        result.append(" ".join(current))

    return result


def _make_chunk(
    doc_id: str,
    page: int,
    idx: int,
    text: str,
    chunk_type: str,
    section: str,
) -> dict:
    return {
        "id": f"{doc_id}_{page:04d}_{idx:04d}",
        "text": text,
        "chunk_type": chunk_type,
        "citation_label": f"p. {page}",
        "citation_page": page,
        "citation_slide": 0,
        "citation_section": section,
        "doc_id": doc_id,
    }


def _make_parent(doc_id: str, parent_idx: int, page: int, text: str, section: str) -> dict:
    return {
        "id": f"{doc_id}_p{parent_idx:04d}",
        "text": text,
        "doc_id": doc_id,
        "citation_label": f"p. {page}",
        "citation_page": page,
        "citation_section": section,
    }


def _make_child(
    parent_id: str,
    doc_id: str,
    page: int,
    child_idx: int,
    text: str,
    section: str,
    chunk_type: str = "paragraph",
) -> dict:
    return {
        "id": f"{parent_id}_c{child_idx:04d}",
        "text": text,
        "chunk_type": chunk_type,
        "parent_id": parent_id,
        "citation_label": f"p. {page}",
        "citation_page": page,
        "citation_slide": 0,
        "citation_section": section,
        "doc_id": doc_id,
    }


def chunk_pages(pages: list[dict], doc_id: str) -> list[dict]:
    """Legacy flat-chunk function. Superseded by build_parent_child_chunks() in run().

    Mantenuta per compatibilità con eventuali chiamate esterne e test.
    """
    chunks: list[dict] = []
    current_section = ""
    chunk_idx = 0

    for page_data in pages:
        page_num = page_data["page"]
        text = page_data["text"]

        if not text.strip():
            continue

        blocks = _split_into_blocks(text)
        pending_heading = ""

        for block in blocks:
            btype = block["type"]
            btext = block["text"].strip()

            if not btext:
                continue

            if btype == "heading":
                heading_text = re.sub(r'^#{1,4}\s+', '', btext).strip()
                current_section = heading_text
                pending_heading = btext
                continue

            if btype == "table":
                full_text = f"{pending_heading}\n\n{btext}" if pending_heading else btext
                pending_heading = ""
                if _word_count(full_text) >= _MIN_WORDS_TABLE:
                    chunks.append(_make_chunk(doc_id, page_num, chunk_idx, full_text, "table", current_section))
                    chunk_idx += 1
                continue

            full_text = f"{pending_heading}\n\n{btext}" if pending_heading else btext
            pending_heading = ""

            for sub in _split_paragraph(full_text, _MAX_WORDS, overlap_words=_OVERLAP_WORDS):
                sub = sub.strip()
                if _word_count(sub) >= _MIN_WORDS:
                    chunks.append(_make_chunk(doc_id, page_num, chunk_idx, sub, "paragraph", current_section))
                    chunk_idx += 1

    return chunks


def _clean_section_title(title: str) -> str:
    """Strip markdown emphasis residue from heading titles (e.g. '**Prologue**')."""
    return re.sub(r'\*{1,3}|_{1,3}|`', '', title or '').strip()


def build_parent_child_chunks(
    pages: list[dict],
    doc_id: str,
) -> tuple[list[dict], list[dict], list[dict]]:
    """Chunking gerarchico: parent (1200 parole) → child (150 parole).

    I child chunk sono le unità di retrieval indicizzate in QVAC.
    I parent chunk forniscono contesto esteso al LLM dopo il retrieval.

    Oltre ai chunk, raccoglie gli eventi sezione in ordine di documento:
    ogni heading markdown produce {"title", "level", "page_start",
    "page_end", "parent_chunk_ids"}; i parent generati prima del primo
    heading finiscono in una sezione preambolo con title="". La lista
    piatta viene annidata da build_section_tree().

    Ritorna (parent_chunks, child_chunks, section_events).
    """
    parents: list[dict] = []
    children: list[dict] = []
    section_events: list[dict] = []
    current_section = ""
    parent_idx = 0

    def _record_parent_in_section(parent: dict) -> None:
        if not section_events:
            # Content before the first heading — untitled preamble section.
            section_events.append({
                "title": "",
                "level": 1,
                "page_start": parent["citation_page"],
                "page_end": parent["citation_page"],
                "parent_chunk_ids": [],
            })
        event = section_events[-1]
        event["parent_chunk_ids"].append(parent["id"])
        event["page_end"] = max(event["page_end"], parent["citation_page"])

    for page_data in pages:
        page_num = page_data["page"]
        text = page_data["text"]

        if not text.strip():
            continue

        blocks = _split_into_blocks(text)
        pending_heading = ""

        for block in blocks:
            btype = block["type"]
            btext = block["text"].strip()

            if not btext:
                continue

            if btype == "heading":
                level = len(re.match(r'^#{1,4}', btext).group(0))
                heading_text = re.sub(r'^#{1,4}\s+', '', btext).strip()
                current_section = heading_text
                pending_heading = btext
                section_events.append({
                    "title": _clean_section_title(heading_text),
                    "level": level,
                    "page_start": page_num,
                    "page_end": page_num,
                    "parent_chunk_ids": [],
                })
                continue

            full_text = f"{pending_heading}\n\n{btext}" if pending_heading else btext
            pending_heading = ""

            # Suddividi in parent block (nessun overlap tra parent)
            for parent_text in _split_paragraph(full_text, _PARENT_WORDS, overlap_words=0):
                parent_text = parent_text.strip()
                if not parent_text:
                    continue

                wc = _word_count(parent_text)
                min_thresh = _MIN_WORDS_TABLE if btype == "table" else _MIN_WORDS
                if wc < min_thresh:
                    continue

                parent = _make_parent(doc_id, parent_idx, page_num, parent_text, current_section)
                parents.append(parent)
                _record_parent_in_section(parent)

                # Genera child chunk dal parent
                if btype in ("table", "formula", "code") or wc <= _CHILD_WORDS:
                    # Tabelle, formule, code block e blocchi piccoli: un solo child = il parent intero
                    child = _make_child(
                        parent["id"], doc_id, page_num, 0,
                        parent_text, current_section, chunk_type=btype,
                    )
                    if _word_count(child["text"]) >= min_thresh:
                        children.append(child)
                else:
                    child_subs = _split_paragraph(
                        parent_text, _CHILD_WORDS, overlap_words=_CHILD_OVERLAP
                    )
                    ci = 0
                    for child_text in child_subs:
                        child_text = child_text.strip()
                        if _word_count(child_text) < _MIN_WORDS:
                            continue
                        # Hard cap: single long sentences can exceed the target;
                        # split further by words so no chunk exceeds 512 tokens.
                        if _word_count(child_text) > _CHILD_MAX_WORDS:
                            words = child_text.split()
                            step = _CHILD_MAX_WORDS - _CHILD_OVERLAP
                            for j in range(0, len(words), step):
                                seg = " ".join(words[j : j + _CHILD_MAX_WORDS])
                                if _word_count(seg) >= _MIN_WORDS:
                                    children.append(
                                        _make_child(
                                            parent["id"], doc_id, page_num, ci,
                                            seg, current_section,
                                        )
                                    )
                                    ci += 1
                        else:
                            children.append(
                                _make_child(
                                    parent["id"], doc_id, page_num, ci,
                                    child_text, current_section,
                                )
                            )
                            ci += 1

                parent_idx += 1

    return parents, children, section_events


def build_section_tree(section_events: list[dict]) -> list[dict]:
    """Annida gli eventi sezione piatti in un albero per livello heading.

    Ogni nodo: {"title", "level", "page_start", "page_end",
    "parent_chunk_ids", "children"}. Un heading di livello L diventa figlio
    dell'ultimo heading aperto con livello < L (stack-based, ordine di
    documento). page_end è propagato bottom-up sui discendenti, così lo span
    di un capitolo copre tutte le sue sottosezioni.

    È la fonte della struttura per il course builder (outline generation):
    parent_chunk_ids àncora ogni sezione ai ChunkParent da cui generare.
    """
    roots: list[dict] = []
    stack: list[dict] = []

    for event in section_events:
        node = {**event, "parent_chunk_ids": list(event["parent_chunk_ids"]), "children": []}
        while stack and stack[-1]["level"] >= node["level"]:
            stack.pop()
        if stack:
            stack[-1]["children"].append(node)
        else:
            roots.append(node)
        stack.append(node)

    def _propagate_page_end(node: dict) -> int:
        for child in node["children"]:
            node["page_end"] = max(node["page_end"], _propagate_page_end(child))
        return node["page_end"]

    for root in roots:
        _propagate_page_end(root)

    return roots


def build_section_events_from_parents(parent_rows: list) -> list[dict]:
    """Backfill: ricostruisce eventi sezione piatti dalle righe ChunkParent.

    Per i documenti ingestionati prima dell'introduzione del section tree il
    file sorgente può non esistere più (run() lo elimina a fine pipeline),
    ma chunk_parent conserva citation_section e citation_page per parent.
    Raggruppando le run consecutive della stessa sezione (righe ordinate per
    id, che codifica l'ordine di documento) si ottiene un albero piatto di
    livello 1 — senza gerarchia heading, ma sufficiente per l'outline.
    """
    events: list[dict] = []
    for row in parent_rows:
        title = _clean_section_title(row.citation_section)
        page = row.citation_page or 0
        if not events or events[-1]["title"] != title:
            events.append({
                "title": title,
                "level": 1,
                "page_start": page,
                "page_end": page,
                "parent_chunk_ids": [],
            })
        event = events[-1]
        event["parent_chunk_ids"].append(row.id)
        event["page_end"] = max(event["page_end"], page)
    return events


# ---------------------------------------------------------------------------
# Stage 4 — Quality filter
# ---------------------------------------------------------------------------

def filter_chunks(chunks: list[dict]) -> list[dict]:
    """Scarta chunk di qualità insufficiente.

    Criteri di scarto:
    - Paragrafi con meno di _MIN_WORDS parole
    - Tabelle con meno di _MIN_WORDS_TABLE parole
    - Chunk dominati da caption figura (> 60% dei caratteri)
    """
    result = []
    for c in chunks:
        text = c["text"]
        is_table = c.get("chunk_type") == "table"
        threshold = _MIN_WORDS_TABLE if is_table else _MIN_WORDS

        if _word_count(text) < threshold:
            continue

        caption_chars = sum(len(m.group()) for m in _FIGURE_CAPTION_RE.finditer(text))
        if len(text) > 0 and caption_chars / len(text) > 0.6:
            continue

        result.append(c)

    return result


# ---------------------------------------------------------------------------
# Contextual chunk enrichment (Anthropic method — RAG_CONTEXTUAL_CHUNKS=true)
# ---------------------------------------------------------------------------

def _enrich_with_context(
    child_chunks: list[dict],
    doc_summary: str,
    qvac_url: str,
) -> list[dict]:
    """Prepend a short AI-generated context sentence to each child chunk before embedding.

    Calls QVAC /generate once per chunk. Each chunk whose enrichment call fails
    is returned unchanged (fail-open: best-effort, never blocks ingest).
    Only active when RAG_CONTEXTUAL_CHUNKS=true.
    """
    _PROMPT = (
        "Scrivi in una frase (max 25 parole) il contesto di questo estratto: "
        "indica l'argomento e la sezione. Usa la stessa lingua del testo."
    )
    enriched: list[dict] = []
    for chunk in child_chunks:
        chunk_text = chunk["text"]
        ctx_input = f"Documento:\n{doc_summary[:400]}\n\nEstratto:\n{chunk_text[:600]}"
        try:
            with httpx.Client(timeout=_CONTEXTUAL_TIMEOUT) as client:
                resp = client.post(
                    f"{qvac_url}/generate",
                    json={"question": _PROMPT, "context": [{"label": "estratto", "text": ctx_input}]},
                )
                resp.raise_for_status()
                ctx_sentence = resp.json().get("answer", "").strip()
        except Exception as exc:
            logger.debug("Contextual enrichment skipped for chunk %s: %s", chunk.get("id"), exc)
            ctx_sentence = ""

        if ctx_sentence and 5 < len(ctx_sentence) < 300:
            enriched.append({**chunk, "text": f"{ctx_sentence}\n\n{chunk_text}"})
        else:
            enriched.append(chunk)

    logger.info(
        "Contextual enrichment done for %d chunks (%d enriched)",
        len(child_chunks),
        sum(1 for o, n in zip(child_chunks, enriched) if o["text"] != n["text"]),
    )
    return enriched


# ---------------------------------------------------------------------------
# BM25 index
# ---------------------------------------------------------------------------

def _build_bm25_index(child_chunks: list[dict], workspace: str, doc_id: str) -> None:
    """Aggiorna il corpus BM25 per il workspace e ricostruisce l'indice su disco.

    Il corpus (corpus.json) accumula i child chunk di tutti i documenti del corso.
    Su re-ingest dello stesso doc_id, i vecchi entry vengono prima rimossi.
    """
    try:
        from rank_bm25 import BM25Okapi  # type: ignore[import]
    except ImportError:
        logger.warning("rank_bm25 non installato — BM25 index non costruito")
        return

    corpus_path = QVAC_INGEST_DIR / f"{workspace}_corpus.json"
    bm25_path = QVAC_INGEST_DIR / f"{workspace}_bm25.pkl"

    corpus: dict[str, dict] = {}
    if corpus_path.exists():
        try:
            with corpus_path.open(encoding="utf-8") as f:
                corpus = json.load(f)
        except (json.JSONDecodeError, OSError):
            corpus = {}

    # Rimuovi entry stale per questo doc_id (gestione re-ingest)
    corpus = {cid: info for cid, info in corpus.items() if info.get("doc_id") != doc_id}

    # Aggiungi nuovi child chunk
    for c in child_chunks:
        corpus[c["id"]] = {
            "text": c["text"],
            "label": c["citation_label"],
            "page": c["citation_page"],
            "section": c["citation_section"],
            "doc_id": c["doc_id"],
            "parent_id": c.get("parent_id", ""),
        }

    with corpus_path.open("w", encoding="utf-8") as f:
        json.dump(corpus, f, ensure_ascii=False)

    from app.services.hybrid_search import _tokenize  # noqa: PLC0415
    ids = list(corpus.keys())
    tokenized = [_tokenize(corpus[i]["text"]) for i in ids]
    bm25 = BM25Okapi(tokenized)

    import pickle
    with bm25_path.open("wb") as f:
        pickle.dump({"ids": ids, "bm25": bm25}, f)

    logger.info("BM25 index aggiornato per workspace '%s': %d chunk totali", workspace, len(ids))


# ---------------------------------------------------------------------------
# Parent DB helpers
# ---------------------------------------------------------------------------

def _save_parents_to_db(parents: list[dict], course_id: str, db) -> None:
    """Salva i parent chunk nella tabella ChunkParent (upsert per id)."""
    from app.db.models import ChunkParent  # noqa: PLC0415

    for p in parents:
        existing = db.query(ChunkParent).filter_by(id=p["id"]).first()
        if existing:
            existing.text = p["text"]
            existing.course_id = course_id
            existing.citation_label = p["citation_label"]
            existing.citation_page = p["citation_page"]
            existing.citation_section = p["citation_section"]
        else:
            db.add(ChunkParent(
                id=p["id"],
                doc_id=p["doc_id"],
                course_id=course_id,
                text=p["text"],
                citation_label=p["citation_label"],
                citation_page=p["citation_page"],
                citation_section=p["citation_section"],
            ))
    db.commit()
    logger.debug("Saved %d parent chunks to DB for course '%s'", len(parents), course_id)


# ---------------------------------------------------------------------------
# JSONL helpers
# ---------------------------------------------------------------------------

def _write_jsonl(chunks: list[dict], document_id: str) -> Path:
    QVAC_INGEST_DIR.mkdir(parents=True, exist_ok=True)
    out_path = QVAC_INGEST_DIR / f"{document_id}_contingency.jsonl"
    with out_path.open("w", encoding="utf-8") as f:
        for chunk in chunks:
            f.write(json.dumps(chunk) + "\n")
    logger.debug("Wrote %d chunks to %s", len(chunks), out_path)
    return out_path


# ---------------------------------------------------------------------------
# QVAC ingest
# ---------------------------------------------------------------------------

def _qvac_ingest(jsonl_path: Path, workspace: str, rebuild: bool = False) -> bool:
    """POST the JSONL path to QVAC for embedding + HyperDB indexing.

    Returns True on success, False on any network/HTTP error.
    Timeout configurable via QVAC_INGEST_TIMEOUT env (default 300s).
    """
    try:
        timeout = float(os.getenv("QVAC_INGEST_TIMEOUT", "300"))
        with httpx.Client(timeout=timeout) as client:
            resp = client.post(
                f"{QVAC_SERVICE_URL}/ingest",
                json={"jsonlPath": str(jsonl_path), "workspace": workspace, "rebuild": rebuild},
            )
            resp.raise_for_status()
            logger.info("QVAC ingest accepted — workspace '%s', status %d", workspace, resp.status_code)
            return True
    except httpx.HTTPError as exc:
        logger.warning("QVAC service unavailable, skipping QVAC ingest: %s", exc)
        return False


# ---------------------------------------------------------------------------
# DB helpers
# ---------------------------------------------------------------------------

def _set_stage(doc: CourseDocument, stage: DocumentProcessingStage, db) -> None:
    doc.processing_stage = stage
    db.commit()


def _mark_error(doc: CourseDocument, message: str, db) -> None:
    doc.status = DocumentStatus.ERROR
    doc.processing_stage = DocumentProcessingStage.ERROR
    doc.error_message = message
    db.commit()


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def run(
    document_id: str,
    course_id: str,
    filename: str,
    file_path: str,
    material_type: str = "lecture",
) -> None:
    """Execute the full ingestion pipeline for an uploaded document."""
    logger.info("Pipeline starting for document %s (%s)", document_id, filename)

    with get_db_context() as db:
        doc = document_repo.get_by_id(db, document_id)
        if doc is None:
            logger.error("Document %s not found — aborting pipeline", document_id)
            return

        ext = Path(filename).suffix.lower()
        if ext not in {".pdf", ".pptx", ".docx"}:
            _mark_error(doc, f"Unsupported file type '{ext}'. Accepted: PDF, PPTX, DOCX.", db)
            return

        try:
            # ------------------------------------------------------------------
            # Stage 1 — PARSING (page-by-page)
            # ------------------------------------------------------------------
            _set_stage(doc, DocumentProcessingStage.PARSING, db)

            if ext == ".pdf":
                if _USE_DOCLING:
                    try:
                        pages, page_count = _parse_pdf_with_docling(file_path)
                        parser_used = "docling"
                    except Exception as exc:
                        logger.warning(
                            "Docling failed for %s (%s) — falling back to pymupdf4llm",
                            document_id, exc,
                        )
                        pages, page_count = parse_pdf_pages(file_path)
                        parser_used = "pymupdf4llm-page-chunks"
                else:
                    pages, page_count = parse_pdf_pages(file_path)
                    parser_used = "pymupdf4llm-page-chunks"
            elif ext == ".pptx":
                pages, page_count = parse_pptx_pages(file_path)
                parser_used = "python-pptx"
            else:
                pages, page_count = parse_docx_pages(file_path)
                parser_used = "python-docx"

            logger.info("Parsed %d pages for %s", len(pages), document_id)

            # ------------------------------------------------------------------
            # Stage 2 — CLEANING (PDF only; PPTX/DOCX are already clean)
            # ------------------------------------------------------------------
            if ext == ".pdf":
                boilerplate = detect_boilerplate(pages)
                if boilerplate:
                    logger.info("Detected %d boilerplate lines for %s", len(boilerplate), document_id)
                pages = [
                    {"page": p["page"], "text": clean_page(p["text"], boilerplate)}
                    for p in pages
                ]

            # ------------------------------------------------------------------
            # Stage 3 — CHUNKING (parent-child hierarchy)
            # ------------------------------------------------------------------
            _set_stage(doc, DocumentProcessingStage.CHUNKING, db)
            parent_chunks, raw_children, section_events = build_parent_child_chunks(
                pages, document_id
            )

            # ------------------------------------------------------------------
            # Stage 3b — QUALITY FILTER (child chunks only)
            # ------------------------------------------------------------------
            chunks = filter_chunks(raw_children)
            dropped = len(raw_children) - len(chunks)
            logger.info(
                "Parent-child for %s: %d parents, %d children raw → %d after filter (%d dropped)",
                document_id, len(parent_chunks), len(raw_children), len(chunks), dropped,
            )

            # ------------------------------------------------------------------
            # Stage 3c — CONTEXTUAL ENRICHMENT (opt-in: RAG_CONTEXTUAL_CHUNKS=true)
            # ------------------------------------------------------------------
            if _RAG_CONTEXTUAL_CHUNKS and chunks:
                doc_summary = "\n".join(p["text"][:200] for p in pages[:3])
                logger.info(
                    "Contextual enrichment enabled — enriching %d chunks for %s",
                    len(chunks), document_id,
                )
                chunks = _enrich_with_context(chunks, doc_summary, QVAC_SERVICE_URL)

            # ------------------------------------------------------------------
            # Stage 3d — SAVE PARENTS TO DB
            # ------------------------------------------------------------------
            try:
                _save_parents_to_db(parent_chunks, course_id, db)
            except Exception as exc:
                logger.warning("Could not save parent chunks to DB for %s: %s", document_id, exc)

            # ------------------------------------------------------------------
            # Stage 5 — INDEXING (ChromaDB — skipped when SKIP_CHROMA_INDEX=true)
            # ------------------------------------------------------------------
            _set_stage(doc, DocumentProcessingStage.INDEXING, db)

            if not SKIP_CHROMA_INDEX:
                from fastembed import TextEmbedding  # noqa: PLC0415
                import chromadb                      # noqa: PLC0415
                from chromadb.config import Settings as ChromaSettings  # noqa: PLC0415

                os.makedirs(CHROMA_DB_PATH, exist_ok=True)
                chroma_client = chromadb.PersistentClient(
                    path=CHROMA_DB_PATH,
                    settings=ChromaSettings(anonymized_telemetry=False),
                )
                collection = chroma_client.get_or_create_collection(
                    name=CHROMA_COLLECTION_NAME,
                    metadata={"hnsw:space": "cosine"},
                )
                embedding_model = TextEmbedding("sentence-transformers/all-MiniLM-L6-v2")
                texts = [c["text"] for c in chunks]
                embeddings = [v.tolist() for v in embedding_model.embed(texts)]
                ids = [c["id"] for c in chunks]
                metadatas = [
                    {
                        "doc_id": c["doc_id"],
                        "filename": filename,
                        "course_id": course_id,
                        "material_type": material_type,
                        "label": c["citation_label"],
                        "section": c["citation_section"],
                        "page": c["citation_page"],
                        "slide": c["citation_slide"],
                        "chunk_type": c["chunk_type"],
                        "parent_id": c.get("parent_id", ""),
                    }
                    for c in chunks
                ]
                if ids:
                    collection.upsert(
                        ids=ids,
                        embeddings=embeddings,
                        documents=texts,
                        metadatas=metadatas,  # type: ignore[arg-type]
                    )
                    logger.info("Indexed %d vectors into ChromaDB at %s", len(ids), CHROMA_DB_PATH)
            else:
                logger.info("SKIP_CHROMA_INDEX=true — skipping ChromaDB embedding for %s", document_id)

            # ------------------------------------------------------------------
            # Stage 6 — QVAC ingest
            # ------------------------------------------------------------------
            jsonl_path = _write_jsonl(chunks, document_id)
            qvac_ok = _qvac_ingest(jsonl_path, workspace=course_id, rebuild=False)

            # ------------------------------------------------------------------
            # Stage 6b — BM25 index update (sparse retrieval)
            # ------------------------------------------------------------------
            try:
                _build_bm25_index(chunks, course_id, document_id)
            except Exception as exc:
                logger.warning("BM25 index build failed for %s: %s", document_id, exc)

            # ------------------------------------------------------------------
            # Finalise DB record
            # ------------------------------------------------------------------
            full_text = "\n\n".join(p["text"] for p in pages)
            sections = list(dict.fromkeys(          # dedup preserving order
                c["citation_section"] for c in chunks if c["citation_section"]
            ))[:20]
            sample = [
                {
                    "text": c["text"][:300],
                    "label": c["citation_label"],
                    "section": c["citation_section"],
                }
                for c in chunks[:5]
            ]

            doc.status = DocumentStatus.READY
            doc.processing_stage = DocumentProcessingStage.DONE
            doc.indexing_status = "indexed" if qvac_ok else "qvac_pending"
            doc.chunk_count = len(chunks)
            doc.parser_used = parser_used
            doc.page_count = page_count if page_count else None
            doc.extracted_text_preview = full_text[:500]
            doc.sections_json = json.dumps(sections)
            doc.section_tree_json = json.dumps(build_section_tree(section_events))
            doc.sample_chunks_json = json.dumps(sample)
            db.commit()

            logger.info(
                "Pipeline done for %s — %d chunks, parser=%s, qvac_ok=%s",
                document_id, len(chunks), parser_used, qvac_ok,
            )

            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except OSError:
                pass

        except Exception as exc:
            logger.exception("Pipeline failed for %s: %s", document_id, exc)
            _mark_error(doc, str(exc), db)


# ---------------------------------------------------------------------------
# QVAC-only reindex — retries the QVAC ingest step without re-parsing
# ---------------------------------------------------------------------------

def reindex_qvac(document_id: str, course_id: str) -> None:
    """Retry QVAC ingest for a document whose indexing_status is 'qvac_pending'."""
    jsonl_path = QVAC_INGEST_DIR / f"{document_id}_contingency.jsonl"
    if not jsonl_path.exists():
        logger.warning("Cannot reindex %s: JSONL not found at %s", document_id, jsonl_path)
        return

    with get_db_context() as db:
        doc = document_repo.get_by_id(db, document_id)
        if doc is None:
            logger.error("Document %s not found — aborting reindex", document_id)
            return

        qvac_ok = _qvac_ingest(jsonl_path, workspace=course_id, rebuild=True)
        doc.indexing_status = "indexed" if qvac_ok else "qvac_pending"
        db.commit()
        logger.info("Reindex QVAC for %s: indexing_status=%s", document_id, doc.indexing_status)
