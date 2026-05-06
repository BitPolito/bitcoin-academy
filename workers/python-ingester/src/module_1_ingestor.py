import os
import logging
import fitz  # PyMuPDF — 4-5x faster than pdfplumber for text+size extraction
from pptx import Presentation

logger = logging.getLogger(__name__)

class RamSafeIngestor:
    def __init__(self, file_path: str, chunk_size: int = 100):
        self.file_path = file_path
        self.chunk_size = chunk_size

        if not os.path.exists(self.file_path):
            raise FileNotFoundError(f"Target file not found at {self.file_path}")

        self.is_pptx = self.file_path.lower().endswith('.pptx')
        self.total_pages = self._calculate_total_pages()

    def _calculate_total_pages(self) -> int:
        logger.info(f"Calculating total pages/slides for: {os.path.basename(self.file_path)}")
        if self.is_pptx:
            prs = Presentation(self.file_path)
            return len(prs.slides)
        else:
            with fitz.open(self.file_path) as doc:
                return len(doc)

    def process_in_batches(self, batch_callback):
        logger.info(f"Target locked: {self.total_pages} pages/slides detected.")

        if self.is_pptx:
            prs = Presentation(self.file_path)
            logger.info("Opening PPTX presentation...")
            batch_callback(prs.slides)
        else:
            # Keep the document open for the full run; fitz.Page objects are only
            # valid while their parent Document is open.
            with fitz.open(self.file_path) as doc:
                for i in range(0, self.total_pages, self.chunk_size):
                    end_idx = min(i + self.chunk_size, self.total_pages)
                    logger.info(f"Opening memory-safe slice: Pages {i + 1} to {end_idx}...")
                    batch = [doc[j] for j in range(i, end_idx)]
                    batch_callback(batch)