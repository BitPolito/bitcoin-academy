import uuid
import statistics
import logging
import re
from typing import Any, List, Optional
from datetime import datetime
from services.ai.app.schemas.normalized_document import (
    NormalizedDocument,
    DocumentBlock,
    BlockPosition,
    BlockType,
    DocumentType,
    TextbookExcerptMetadata
)

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docling.document_converter import DocumentConverter, PdfFormatOption
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions, EasyOcrOptions
except ImportError:
    DocumentConverter = None
    PdfFormatOption = None
    InputFormat = None
    PdfPipelineOptions = None
    EasyOcrOptions = None

try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO_SHAPE_TYPE
except ImportError:
    _MSO_SHAPE_TYPE = None

logger = logging.getLogger(__name__)

_ocr_reader = None

def _get_ocr_reader():
    global _ocr_reader
    if _ocr_reader is None:
        import easyocr
        _ocr_reader = easyocr.Reader(["en", "it"], gpu=False, verbose=False)
    return _ocr_reader

def _ocr_image_blob(blob: bytes) -> str:
    import numpy as np
    from PIL import Image
    import io
    try:
        img = Image.open(io.BytesIO(blob)).convert("RGB")
        arr = np.array(img)
        results = _get_ocr_reader().readtext(arr, detail=0, paragraph=True)
        return " ".join(str(t).strip() for t in results if str(t).strip())
    except Exception as exc:
        logger.debug("OCR on image blob failed: %s", exc)
        return ""

class StructuralParser:
    def __init__(
        self,
        course_id: str,
        document_id: str,
        document_type: DocumentType,
        title: str,
        source_filename: str,
        file_path: str,
        use_advanced_parser: bool,
        lecture_id: Optional[str] = None,
        tags: Optional[List[str]] = None,
        prerequisites: Optional[List[str]] = None,
    ):
        self.course_id = course_id
        self.document_id = document_id
        self.document_type = document_type
        self.title = title
        self.source_filename = source_filename
        self.file_path = file_path
        self.use_advanced_parser = use_advanced_parser
        self.lecture_id = lecture_id or document_id
        self.tags = tags or []
        self.prerequisites = prerequisites or []
        self.current_section_path = []
        self.in_exclusion_zone = False

        
    def _sanitize_text(self, text: str) -> str:
        # Fix PDF ligature corruption (the smoking gun from Query 20)
        text = text.replace('昀椀', 'fi').replace('昀氀', 'fl')
        # Fix squished camelCase headers (e.g., "DefiningTheLoss" -> "Defining The Loss")
        text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
        return text
    
    def _parse_with_docling(self) -> NormalizedDocument:
        if not DocumentConverter or not self.file_path:
            raise ValueError("Docling is not installed or file_path is missing")

        from docling_core.types.doc.labels import DocItemLabel
        from docling_core.types.doc.document import TableItem, DocItem

        LABEL_MAP = {
            DocItemLabel.TITLE:          (BlockType.HEADING,    1),
            DocItemLabel.SECTION_HEADER: (BlockType.HEADING,    2),
            DocItemLabel.TEXT:           (BlockType.PARAGRAPH,  None),
            DocItemLabel.PARAGRAPH:      (BlockType.PARAGRAPH,  None),
            DocItemLabel.LIST_ITEM:      (BlockType.LIST_ITEM,  None),
            DocItemLabel.TABLE:          (BlockType.TABLE,      None),
            DocItemLabel.FORMULA:        (BlockType.MATH,       None),
            DocItemLabel.CODE:           (BlockType.CODE_BLOCK, None),
        }

        format_options = {}
        if PdfFormatOption and PdfPipelineOptions and EasyOcrOptions and InputFormat:
            pdf_opts = PdfPipelineOptions(
                do_ocr=True,
                ocr_options=EasyOcrOptions(
                    lang=["en", "it"],
                    force_full_page_ocr=True,
                    confidence_threshold=0.4,
                ),
            )
            format_options[InputFormat.PDF] = PdfFormatOption(pipeline_options=pdf_opts)

        converter = DocumentConverter(format_options=format_options or None)
        doc = converter.convert(self.file_path).document
        blocks = []

        for node, _ in doc.iterate_items():
            # iterate_items yields NodeItem; only DocItems carry label/prov/text
            if not isinstance(node, DocItem):
                continue
            item: DocItem = node

            block_type, default_heading_level = LABEL_MAP.get(
                item.label, (BlockType.PARAGRAPH, None)
            )

            page_num = item.prov[0].page_no if item.prov else None

            if isinstance(item, TableItem):
                text = item.export_to_markdown()
            else:
                raw = getattr(item, "text", None)
                if not raw:
                    continue
                text = self._sanitize_text(str(raw))

            if not text.strip():
                continue

            heading_level = None
            if block_type == BlockType.HEADING:
                heading_level = getattr(item, "level", default_heading_level) or default_heading_level or 1
                self.current_section_path = [text]

            blocks.append(DocumentBlock(
                block_id=str(uuid.uuid4()),
                block_type=block_type,
                text=text,
                position=BlockPosition(
                    page=page_num,
                    section_path=self.current_section_path.copy(),
                ),
                heading_level=heading_level,
            ))

        page_count = len({b.position.page for b in blocks if b.position.page is not None}) or None

        # Supplemental OCR pass: extract text from PICTURE shapes that Docling's
        # MsPowerpointDocumentBackend silently discards (no image pipeline for PPTX).
        if self.file_path and self.file_path.lower().endswith((".pptx", ".ppt")):
            try:
                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
                prs = Presentation(self.file_path)
                for slide_num, slide in enumerate(prs.slides, 1):
                    for shape in slide.shapes:
                        if getattr(shape, "shape_type", None) not in (
                            _MSO.PICTURE, _MSO.LINKED_PICTURE
                        ):
                            continue
                        try:
                            ocr_text = self._sanitize_text(_ocr_image_blob(shape.image.blob))
                            if ocr_text:
                                blocks.append(DocumentBlock(
                                    block_id=str(uuid.uuid4()),
                                    block_type=BlockType.PARAGRAPH,
                                    text=ocr_text,
                                    position=BlockPosition(
                                        page=slide_num,
                                        section_path=[f"Slide {slide_num}"],
                                    ),
                                ))
                        except Exception as exc:
                            logger.debug(
                                "OCR failed for image shape on slide %d: %s", slide_num, exc
                            )
            except Exception as exc:
                logger.warning("Supplemental PPTX OCR pass failed: %s", exc)

        return NormalizedDocument(
            doc_id=self.document_id,
            course_id=self.course_id,
            lecture_id=self.lecture_id,
            document_type=self.document_type,
            title=self.title,
            source_filename=self.source_filename,
            parser_used="docling-tier-2-parser",
            parsed_at=datetime.now(),
            page_count=page_count,
            slide_count=None,
            blocks=blocks,
            tags=self.tags,
            prerequisites=self.prerequisites,
            type_metadata=TextbookExcerptMetadata(book_title=self.title).model_dump() if self.document_type == DocumentType.TEXTBOOK_EXCERPT else None
        )

    def parse_pages(self, pages: List[Any], total_pages: int) -> NormalizedDocument:
        if self.use_advanced_parser:
            try:
                return self._parse_with_docling()
            except Exception as e:
                available_pages = len(pages) if pages is not None else 0
                logger.warning(
                    "Docling failed (%s), falling back to hybrid parser with %d available pages.",
                    e,
                    available_pages,
                )
                if not available_pages:
                    logger.error("Fallback aborted: ingestor provided 0 pages — re-raising Docling exception.")
                    raise

        blocks = []
        
        # --- PPTX HANDLING ---
        if self.document_type == DocumentType.LECTURE_SLIDES:
            for i, slide in enumerate(pages):
                slide_num = i + 1
                title_text = slide.shapes.title.text.strip() if slide.shapes.title else f"Slide {slide_num}"
                self.current_section_path = [self._sanitize_text(title_text)]

                blocks.append(DocumentBlock(
                    block_id=str(uuid.uuid4()),
                    block_type=BlockType.SLIDE_TITLE,
                    text=self.current_section_path[0],
                    position=BlockPosition(slide=slide_num, section_path=self.current_section_path.copy()),
                    heading_level=1
                ))

                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        body_text = self._sanitize_text(shape.text.strip())
                        if body_text:
                            blocks.append(DocumentBlock(
                                block_id=str(uuid.uuid4()),
                                block_type=BlockType.SLIDE_BODY,
                                text=body_text,
                                position=BlockPosition(slide=slide_num, section_path=self.current_section_path.copy())
                            ))
                    elif _MSO_SHAPE_TYPE and getattr(shape, "shape_type", None) in (
                        _MSO_SHAPE_TYPE.PICTURE,
                        _MSO_SHAPE_TYPE.LINKED_PICTURE,
                    ):
                        try:
                            ocr_text = self._sanitize_text(_ocr_image_blob(shape.image.blob))
                            if ocr_text:
                                blocks.append(DocumentBlock(
                                    block_id=str(uuid.uuid4()),
                                    block_type=BlockType.SLIDE_BODY,
                                    text=ocr_text,
                                    position=BlockPosition(slide=slide_num, section_path=self.current_section_path.copy())
                                ))
                        except Exception as exc:
                            logger.debug("Skipping image shape on slide %d: %s", slide_num, exc)

                # Speaker notes often contain the actual explanations
                try:
                    notes_text = self._sanitize_text(
                        slide.notes_slide.notes_text_frame.text.strip()
                    )
                    if notes_text:
                        blocks.append(DocumentBlock(
                            block_id=str(uuid.uuid4()),
                            block_type=BlockType.SPEAKER_NOTES,
                            text=notes_text,
                            position=BlockPosition(slide=slide_num, section_path=self.current_section_path.copy())
                        ))
                except Exception:
                    pass
                        
        # --- PDF HANDLING ---
        else:
            for page in pages:
                # Reset exclusion zone at each new page so content after the last
                # "References" heading on a previous page is not silently discarded.
                self.in_exclusion_zone = False

                # fitz.Page: number is 0-indexed; add 1 for human-readable page numbers.
                page_number = page.number + 1

                # Extract word-like objects from fitz spans, preserving size per word.
                words = []
                for block in page.get_text("dict")["blocks"]:
                    if block.get("type") != 0:  # skip image blocks
                        continue
                    for line in block["lines"]:
                        for span in line["spans"]:
                            span_text = span["text"]
                            span_size = span["size"]
                            span_top = span["bbox"][1]
                            for token in span_text.split():
                                if token:
                                    words.append({"text": token, "top": span_top, "size": span_size})

                if not words:
                    continue

                sizes = [w['size'] for w in words]
                if not sizes:
                    continue
                median_size = statistics.median(sizes)

                lines_dict: dict = {}
                for word in words:
                    line_y = round(word['top'] / 2) * 2
                    if line_y not in lines_dict:
                        lines_dict[line_y] = []
                    lines_dict[line_y].append(word)

                # Adaptive threshold: lower to ×1.05 when no line on this page
                # exceeds ×1.15 (uniform-font PDFs, scans, slides exported as PDF).
                any_heading_candidate = any(
                    statistics.mean([w['size'] for w in lines_dict[y]]) > median_size * 1.15
                    for y in lines_dict
                    if lines_dict[y]
                )
                heading_threshold = median_size * (1.15 if any_heading_candidate else 1.05)

                sorted_y_coords = sorted(lines_dict.keys())
                content_buffer = []

                for y_coord in sorted_y_coords:
                    line_words = lines_dict[y_coord]
                    raw_line = " ".join(w['text'] for w in line_words).strip()
                    
                    if not raw_line:
                        continue
                        
                    # 1. KILL THE TOC TRAP: Drop lines with leader dots (e.g., ".... 142")
                    if re.search(r'\.{5,}\s*\d+$', raw_line):
                        continue
                        
                    line_text = self._sanitize_text(raw_line)
                    line_size = statistics.mean([w['size'] for w in line_words])

                    # HEADING DETECTION
                    if line_size > heading_threshold and len(line_text.split()) < 15:
                        # 2. THE EXCLUSION ZONE: Check if we just hit an Exercises/References section
                        lower_head = line_text.lower()
                        if "exercises" in lower_head or "references" in lower_head or "bibliography" in lower_head:
                            self.in_exclusion_zone = True
                            content_buffer = [] # Dump whatever we had
                            continue
                        else:
                            self.in_exclusion_zone = False # Back to normal text
                            
                        if content_buffer and not self.in_exclusion_zone:
                            blocks.append(self._build_block(content_buffer, page_number))
                            content_buffer = []

                        self.current_section_path = [line_text]
                        
                        if not self.in_exclusion_zone:
                            blocks.append(DocumentBlock(
                                block_id=str(uuid.uuid4()),
                                block_type=BlockType.HEADING,
                                text=line_text,
                                position=BlockPosition(page=page_number, section_path=self.current_section_path.copy()),
                                heading_level=1
                            ))
                    else:
                        if not self.in_exclusion_zone:
                            content_buffer.append(line_text)

                if content_buffer and not self.in_exclusion_zone:
                    blocks.append(self._build_block(content_buffer, page_number))

        doc = NormalizedDocument(
            doc_id=self.document_id,
            course_id=self.course_id,
            lecture_id=self.lecture_id,
            document_type=self.document_type,
            title=self.title,
            source_filename=self.source_filename,
            parser_used="verilocal-hybrid-parser-v2",
            parsed_at=datetime.now(),
            page_count=total_pages if self.document_type != DocumentType.LECTURE_SLIDES else None,
            slide_count=total_pages if self.document_type == DocumentType.LECTURE_SLIDES else None,
            blocks=blocks,
            tags=self.tags,
            prerequisites=self.prerequisites,
            type_metadata=TextbookExcerptMetadata(book_title=self.title).model_dump() if self.document_type == DocumentType.TEXTBOOK_EXCERPT else None
        )
        return doc

    def _build_block(self, content_buffer: List[str], page_number: int) -> DocumentBlock:
        full_text = " ".join(content_buffer)
        return DocumentBlock(
            block_id=str(uuid.uuid4()),
            block_type=BlockType.PARAGRAPH,
            text=full_text,
            position=BlockPosition(
                page=page_number,
                section_path=self.current_section_path.copy()
            )
        )